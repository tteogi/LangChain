from typing import List
from langchain_community.document_loaders import PyPDFLoader, UnstructuredMarkdownLoader
from langchain.schema import Document
from pptx import Presentation
import tempfile
import os
import re


def parse_faq_markdown(text: str, source_name: str) -> List[Document]:
    """
    FAQ 형식의 마크다운을 파싱하여 각 FAQ 항목을 개별 Document로 변환

    형식: [카테고리] 질문
          답변...
    """
    documents = []

    # [카테고리] 패턴으로 FAQ 항목 분리
    # 패턴: [카테고리] 질문 형태를 찾음
    faq_pattern = r'\[([^\]]+)\]\s*([^\n]+)\n((?:(?!\[)[^\n]*\n?)*)'

    matches = re.finditer(faq_pattern, text, re.MULTILINE)

    for match in matches:
        category = match.group(1).strip()
        question = match.group(2).strip()
        answer = match.group(3).strip()

        # 빈 답변 스킵
        if not answer:
            continue

        # FAQ 전체 내용
        full_content = f"[{category}] {question}\n{answer}"

        # 메타데이터 기본값
        metadata = {
            "source": source_name,
            "category": category,
            "question": question,
            "type": "faq"
        }

        # 수수료 정보 추출 (질문 + 답변에서 "수수료" 키워드가 있는 경우)
        if "수수료" in question or "수수료" in answer:
            # 비율 수수료 추출 (예: 3%, 5.5%)
            percent_match = re.search(r'(\d+(?:\.\d+)?)\s*%', answer)
            if percent_match:
                fee_value = float(percent_match.group(1))
                metadata["fee_type"] = "percent"
                metadata["fee_value"] = fee_value
                metadata["fee_display"] = f"{percent_match.group(1)}%"
            else:
                # 고정 금액 수수료 추출 (예: 1,000원, 1000원)
                won_match = re.search(r'(\d+,?\d*)\s*원', answer)
                if won_match:
                    fee_str = won_match.group(1).replace(',', '')
                    fee_value = float(fee_str)
                    metadata["fee_type"] = "fixed"
                    metadata["fee_value"] = fee_value
                    metadata["fee_display"] = f"{won_match.group(1)}원"

        doc = Document(
            page_content=full_content,
            metadata=metadata
        )
        documents.append(doc)

    return documents


def load_documents(uploaded_files) -> List[Document]:
    documents = []
    
    for uploaded_file in uploaded_files:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            if file_extension == '.pdf':
                loader = PyPDFLoader(tmp_file_path)
                loaded_docs = loader.load()
                for doc in loaded_docs:
                    doc.metadata["source"] = uploaded_file.name
                documents.extend(loaded_docs)
            elif file_extension in ['.md', '.markdown']:
                # FAQ 형식인지 확인
                with open(tmp_file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # [카테고리] 패턴이 여러 개 있으면 FAQ 형식으로 판단
                faq_pattern_count = len(re.findall(r'\[([^\]]+)\]\s*[^\n]+', content))

                if faq_pattern_count >= 3:  # 3개 이상의 FAQ 항목이 있으면
                    # FAQ 구조화 파싱 사용
                    faq_docs = parse_faq_markdown(content, uploaded_file.name)
                    documents.extend(faq_docs)
                else:
                    # 일반 Markdown으로 처리
                    loader = UnstructuredMarkdownLoader(tmp_file_path)
                    loaded_docs = loader.load()
                    for doc in loaded_docs:
                        doc.metadata["source"] = uploaded_file.name
                    documents.extend(loaded_docs)
            elif file_extension == '.pptx':
                prs = Presentation(tmp_file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += shape.text + "\n"
                    
                    if slide_text.strip():
                        doc = Document(
                            page_content=slide_text.strip(),
                            metadata={
                                "source": uploaded_file.name,
                                "page": slide_num - 1,
                                "slide": slide_num
                            }
                        )
                        documents.append(doc)
        finally:
            os.unlink(tmp_file_path)
    
    return documents
